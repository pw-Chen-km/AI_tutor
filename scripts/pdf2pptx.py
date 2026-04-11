#!/usr/bin/env python3
"""
PDF to PPTX conversion script using PyMuPDF and python-pptx directly.
This is more stable than pdf2slides which has paddleocr compatibility issues.
"""
import sys
import json
import base64
import os
import tempfile
import io

def convert_pdf_to_pptx_with_notes(pdf_base64: str, notes_data: list, output_path: str) -> dict:
    """
    Convert PDF to PPTX with speaker notes using PyMuPDF and python-pptx.
    
    Args:
        pdf_base64: Base64 encoded PDF data
        notes_data: List of dicts with {slide_number, note_text}
        output_path: Path to save the output PPTX
    
    Returns:
        dict with success status and message
    """
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from PIL import Image
        
        # Decode PDF
        pdf_bytes = base64.b64decode(pdf_base64)
        
        # Open PDF with PyMuPDF
        pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        # Create new presentation
        prs = Presentation()
        # Set slide dimensions to 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Blank layout
        blank_layout = prs.slide_layouts[6]  # Usually blank layout
        
        # Create notes mapping
        notes_map = {}
        for note in notes_data:
            slide_num = note.get('slide_number', 0)
            note_text = note.get('note_text', '')
            if slide_num and note_text:
                notes_map[slide_num] = note_text
        
        # Process each page
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            
            # Render page to image at high resolution
            # Use a matrix for higher quality (2x zoom)
            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat)
            
            # Convert to PNG bytes
            img_bytes = pix.tobytes("png")
            
            # Create slide
            slide = prs.slides.add_slide(blank_layout)
            
            # Add image to slide (centered, full-width)
            # Calculate dimensions to fit while maintaining aspect ratio
            img_stream = io.BytesIO(img_bytes)
            
            # Get image dimensions
            img = Image.open(img_stream)
            img_width, img_height = img.size
            img_stream.seek(0)  # Reset stream position
            
            # Calculate scaling to fit slide
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Scale to fit width, maintaining aspect ratio
            scale = float(slide_width) / img_width
            new_width = slide_width
            new_height = int(img_height * scale)
            
            # If height exceeds slide height, scale by height instead
            if new_height > slide_height:
                scale = float(slide_height) / img_height
                new_height = slide_height
                new_width = int(img_width * scale)
            
            # Center the image on the slide
            left = (slide_width - new_width) / 2
            top = (slide_height - new_height) / 2
            
            # Add the image
            slide.shapes.add_picture(img_stream, left, top, width=new_width, height=new_height)
            
            # Add notes if available (page_num is 0-indexed, slide_number is 1-indexed)
            slide_number = page_num + 1
            if slide_number in notes_map:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = notes_map[slide_number]
        
        pdf_doc.close()
        
        # Save presentation
        prs.save(output_path)
        
        return {"success": True, "message": f"Successfully converted {len(prs.slides)} slides", "output": output_path}
        
    except ImportError as e:
        return {"success": False, "message": f"Missing dependency: {e}. Please install: pip install PyMuPDF python-pptx Pillow"}
    except Exception as e:
        import traceback
        return {"success": False, "message": f"{str(e)}\n{traceback.format_exc()}"}

def main():
    """Main entry point for CLI usage."""
    if len(sys.argv) < 2:
        print(json.dumps({"success": False, "message": "Missing input JSON"}))
        sys.exit(1)
    
    try:
        # Read input from command line argument (JSON string or file path)
        input_arg = sys.argv[1]
        
        if os.path.exists(input_arg):
            with open(input_arg, 'r', encoding='utf-8') as f:
                input_data = json.load(f)
        else:
            input_data = json.loads(input_arg)
        
        pdf_base64 = input_data.get('pdf_base64', '')
        notes_data = input_data.get('notes', [])
        output_path = input_data.get('output_path', '')
        
        if not pdf_base64:
            print(json.dumps({"success": False, "message": "Missing pdf_base64"}))
            sys.exit(1)
        
        if not output_path:
            # Use temp file if no output path specified
            output_path = tempfile.mktemp(suffix='.pptx')
        
        result = convert_pdf_to_pptx_with_notes(pdf_base64, notes_data, output_path)
        
        # If successful, include the output file as base64
        if result.get('success') and os.path.exists(output_path):
            with open(output_path, 'rb') as f:
                result['pptx_base64'] = base64.b64encode(f.read()).decode('utf-8')
            
            # Clean up output file if it was a temp file
            if 'output_path' not in input_data:
                os.unlink(output_path)
        
        print(json.dumps(result))
        
    except json.JSONDecodeError as e:
        print(json.dumps({"success": False, "message": f"Invalid JSON: {e}"}))
        sys.exit(1)
    except Exception as e:
        import traceback
        print(json.dumps({"success": False, "message": f"{str(e)}\n{traceback.format_exc()}"}))
        sys.exit(1)

if __name__ == '__main__':
    main()
